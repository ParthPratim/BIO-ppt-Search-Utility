from pptx import Presentation 
import subprocess as sp
print("Enter Text to be Searched")
search = input()
print("Search Results ...\n\n")
output = sp.getoutput('ls ppts/*.pptx')
pptx_ls = output.split('\n')
p_buffer=[]
p_buffer.append(("File Name"," Slide Numbers"))
max_len = 0
for ppt_file in pptx_ls:
    prs = Presentation(ppt_file)
    num=0
    sld_ct = []
    for slide in prs.slides:
        num += 1
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            tokens=text_frame.text.split()
            kill = False
            for keyword in search.split():
                for token in tokens:
                    if token.lower().find(keyword.lower()) != -1:
                        sld_ct.append(num)
                        kill=True
                        break
                if kill:
                    break
            if kill:
                break
    out_slds = ""
    for slide_no in sld_ct:
            out_slds += " | " + str(slide_no)
    if len(sld_ct) > 0:
        p_buffer.append((ppt_file,out_slds))
        max_len = max(max_len,len(ppt_file))

d = 0
for buff in p_buffer:
    print(buff[0],end="")
    spaces = max_len - len(buff[0]) 
    for space in range(0,spaces):
        print(" ",end="")
    print("\t"+buff[1])
    if d == 0:
        d = 1
        print("--------------------------------------------------------------")

